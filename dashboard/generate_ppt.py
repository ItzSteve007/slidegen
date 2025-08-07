from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import os
import random

def generate_ppt(theme, image_paths, output_path):
    prs = Presentation()

    # --- Title Slide ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = f"{theme.capitalize()} Presentation"
    slide.placeholders[1].text = "Generated using AI and uploaded content"

    # --- Index Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Index"
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    for idx, img_path in enumerate(image_paths):
        tf.add_paragraph().text = f"Slide {idx+1}: {theme} Image {idx+1}"

    # --- Content Slides for Each Image ---
    for idx, image_path in enumerate(image_paths):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add image
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(8))

        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        tf = title_box.text_frame
        tf.text = f"{theme.capitalize()} - Image {idx+1}"
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True

        # (Optional) Add AI-generated description if needed
        # You can replace the line below with actual AI captioning
        caption = f"This image showcases a part of the {theme} theme."
        caption_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
        cap_tf = caption_box.text_frame
        cap_tf.text = caption
        cap_tf.paragraphs[0].font.size = Pt(18)

    # --- Final Slide ---
    final_slide = prs.slides.add_slide(prs.slide_layouts[0])
    final_slide.shapes.title.text = "Thank You!"
    final_slide.placeholders[1].text = "Hope you enjoyed the presentation."

    # Save the presentation
    prs.save(output_path)
